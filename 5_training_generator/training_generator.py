"""
AI Training Material Generator
--------------------------------
Takes a technical topic and generates a PowerPoint with slide outline,
quiz questions, and a glossary — all powered by Claude AI.
Usage: python training_generator.py --topic "Python for Data Analysis"
   or: streamlit run training_generator.py  (for web UI mode)
"""

import argparse
import json
import os
import sys
import re
from anthropic import Anthropic
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

load_dotenv()

# ── COLOUR PALETTE ────────────────────────────────────────────────────────────

DARK_BLUE  = RGBColor(0x1A, 0x3A, 0x5C)
MID_BLUE   = RGBColor(0x4A, 0x7A, 0xAA)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0x2A, 0x7A, 0x4B)

# ── AI CONTENT GENERATION ─────────────────────────────────────────────────────

def generate_content(topic: str) -> dict:
    """Call Claude to generate structured training content as JSON."""
    client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

    prompt = f"""You are a professional training content designer.
Create a structured training module on the topic: "{topic}"

Return ONLY valid JSON — no markdown, no explanation, no preamble.

Use this exact structure:
{{
  "title": "Full training module title",
  "tagline": "One line description",
  "objectives": ["objective 1", "objective 2", "objective 3", "objective 4"],
  "slides": [
    {{
      "title": "Slide title",
      "bullets": ["bullet 1", "bullet 2", "bullet 3"]
    }}
  ],
  "quiz": [
    {{
      "question": "Question text?",
      "options": ["A) ...", "B) ...", "C) ...", "D) ..."],
      "answer": "A"
    }}
  ],
  "glossary": [
    {{"term": "Term", "definition": "Definition"}}
  ]
}}

Requirements:
- 6 content slides (not counting title/objectives/quiz/glossary slides)
- 5 quiz questions with 4 options each
- 8 glossary terms
- Keep bullets concise (max 12 words each)
- Use professional but accessible language"""

    print(f"[AI] Generating content for: {topic}")
    message = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=2000,
        messages=[{"role": "user", "content": prompt}]
    )

    raw = message.content[0].text.strip()

    # Clean any accidental markdown fences
    raw = re.sub(r"^```json\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)

    data = json.loads(raw)
    print(f"[AI] Generated {len(data['slides'])} slides, {len(data['quiz'])} quiz questions, {len(data['glossary'])} glossary terms.")
    return data


# ── POWERPOINT BUILDER ────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def build_pptx(data: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, DARK_BLUE)

    # Accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(5.8), Inches(13.33), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BLUE
    shape.line.fill.background()

    add_text_box(slide, data["title"],    0.6, 2.2, 12, 1.5, font_size=36, bold=True,  align=PP_ALIGN.CENTER)
    add_text_box(slide, data["tagline"],  0.6, 4.0, 12, 0.8, font_size=18, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "AI-Generated Training Module", 0.6, 6.2, 12, 0.5, font_size=12, color=MID_BLUE, align=PP_ALIGN.CENTER)

    # ── Slide 2: Learning Objectives ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    add_text_box(slide, "Learning Objectives", 0.4, 0.15, 12, 0.9, font_size=28, bold=True)

    for i, obj in enumerate(data["objectives"]):
        add_text_box(slide, f"✓  {obj}", 0.8, 1.5 + i * 1.1, 11.5, 0.9,
                     font_size=16, color=DARK_BLUE)

    # ── Slides 3+: Content Slides ─────────────────────────────────────────────
    for slide_data in data["slides"]:
        slide = prs.slides.add_slide(blank_layout)
        set_bg(slide, WHITE)

        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK_BLUE
        bar.line.fill.background()

        add_text_box(slide, slide_data["title"], 0.4, 0.15, 12, 0.9, font_size=26, bold=True)

        for j, bullet in enumerate(slide_data["bullets"]):
            add_text_box(slide, f"▸  {bullet}", 0.8, 1.5 + j * 1.1, 11.5, 0.9,
                         font_size=15, color=DARK_BLUE)

    # ── Quiz Slide ────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, DARK_BLUE)
    add_text_box(slide, "📝  Knowledge Check", 0.5, 0.2, 12, 1.0, font_size=32, bold=True, align=PP_ALIGN.CENTER)

    for k, q in enumerate(data["quiz"][:3]):  # Show first 3 on one slide
        y = 1.4 + k * 1.8
        add_text_box(slide, f"Q{k+1}: {q['question']}", 0.5, y, 12, 0.6, font_size=14, bold=True, color=LIGHT_BLUE)
        options_text = "   ".join(q["options"])
        add_text_box(slide, options_text, 0.5, y + 0.65, 12, 0.6, font_size=12, color=WHITE)

    # ── Glossary Slide ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_text_box(slide, "📖  Glossary", 0.4, 0.15, 12, 0.9, font_size=28, bold=True)

    # Two columns
    mid = len(data["glossary"]) // 2
    for idx, item in enumerate(data["glossary"][:mid]):
        add_text_box(slide, f"{item['term']}:", 0.5, 1.4 + idx * 0.75, 3.0, 0.6,
                     font_size=13, bold=True, color=DARK_BLUE)
        add_text_box(slide, item["definition"], 3.5, 1.4 + idx * 0.75, 3.0, 0.6,
                     font_size=12, color=DARK_BLUE)
    for idx, item in enumerate(data["glossary"][mid:]):
        add_text_box(slide, f"{item['term']}:", 7.0, 1.4 + idx * 0.75, 3.0, 0.6,
                     font_size=13, bold=True, color=DARK_BLUE)
        add_text_box(slide, item["definition"], 10.0, 1.4 + idx * 0.75, 3.0, 0.6,
                     font_size=12, color=DARK_BLUE)

    prs.save(output_path)
    print(f"[EXPORT] PowerPoint saved: {output_path}")


# ── STREAMLIT UI (optional) ───────────────────────────────────────────────────

def run_streamlit():
    import streamlit as st
    from io import BytesIO

    st.set_page_config(page_title="AI Training Generator", page_icon="🎓", layout="centered")
    st.title("🎓 AI Training Material Generator")
    st.caption("Enter a topic and get a ready-made PowerPoint training deck powered by Claude AI.")

    topic = st.text_input("Training topic", placeholder="e.g. Python for Data Analysis, ESP Design, Machine Learning Basics")
    generate_btn = st.button("⚡ Generate Training Materials", type="primary")

    if generate_btn and topic:
        if not os.getenv("ANTHROPIC_API_KEY"):
            st.error("ANTHROPIC_API_KEY not set in .env file.")
            st.stop()

        with st.spinner("Claude is generating your training content..."):
            try:
                data = generate_content(topic)
            except Exception as e:
                st.error(f"AI generation failed: {e}")
                st.stop()

        st.success("✅ Content generated!")

        # Preview
        st.subheader(f"📋 {data['title']}")
        st.caption(data["tagline"])

        col1, col2 = st.columns(2)
        with col1:
            st.markdown("**Learning Objectives**")
            for obj in data["objectives"]:
                st.markdown(f"✓ {obj}")
        with col2:
            st.markdown("**Slides Generated**")
            for s in data["slides"]:
                st.markdown(f"▸ {s['title']}")

        with st.expander("📝 Quiz Questions"):
            for i, q in enumerate(data["quiz"], 1):
                st.markdown(f"**Q{i}: {q['question']}**")
                for opt in q["options"]:
                    st.markdown(f"  {opt}")
                st.caption(f"Answer: {q['answer']}")

        # Build and offer download
        with st.spinner("Building PowerPoint..."):
            buffer = BytesIO()
            tmp_path = "/tmp/training_output.pptx"
            build_pptx(data, tmp_path)
            with open(tmp_path, "rb") as f:
                buffer.write(f.read())

        st.download_button(
            "⬇️ Download PowerPoint",
            data=buffer.getvalue(),
            file_name=f"{topic.replace(' ', '_')}_training.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    elif generate_btn:
        st.warning("Please enter a topic first.")


# ── ENTRY POINT ───────────────────────────────────────────────────────────────

if "streamlit" in sys.modules and hasattr(sys.modules["streamlit"], "secrets"):
    # Running inside Streamlit
    run_streamlit()
else:
    # CLI mode
    parser = argparse.ArgumentParser(description="AI Training Material Generator")
    parser.add_argument("--topic",  required=True, help="Topic to generate training for")
    parser.add_argument("--output", default="training_output.pptx", help="Output .pptx filename")
    args = parser.parse_args()

    if not os.getenv("ANTHROPIC_API_KEY"):
        print("[ERROR] ANTHROPIC_API_KEY not set. Add it to a .env file.")
        sys.exit(1)

    content = generate_content(args.topic)
    build_pptx(content, args.output)
    print(f"\n✅ Done! Open: {args.output}")
